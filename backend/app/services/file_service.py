import base64
import io
import asyncio
from typing import List, Dict, Any, Callable, AsyncGenerator
import fitz  # PyMuPDF
import pandas as pd
from docx import Document
from pptx import Presentation
from app.core.llm import get_llm
from langchain_core.messages import SystemMessage, HumanMessage
from app.core.logger import logger
from app.core.llm import get_time_instructions

class FileParsingService:
    @staticmethod
    async def parse_file(filename: str, base64_data: str) -> str:
        """Parses various file types and returns their text content."""
        # Remove data URI header if present
        if "," in base64_data:
            base64_data = base64_data.split(",")[1]
        
        file_bytes = base64.b64decode(base64_data)
        file_io = io.BytesIO(file_bytes)
        
        ext = filename.split(".")[-1].lower()
        
        try:
            if ext == "pdf":
                doc = fitz.open(stream=file_bytes, filetype="pdf")
                text = ""
                for page in doc:
                    text += page.get_text() + "\n"
                doc.close()
                return text
            
            elif ext in ["xlsx", "xls"]:
                df = pd.read_excel(file_io)
                return df.to_string()
            
            elif ext == "docx":
                doc = Document(file_io)
                return "\n".join([p.text for p in doc.paragraphs])
            
            elif ext == "pptx":
                prs = Presentation(file_io)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            
            elif ext in ["md", "txt"]:
                return file_bytes.decode("utf-8")
            
            else:
                return f"[Unsupported file type: {ext}]"
        except Exception as e:
            logger.error(f"Error parsing file {filename}: {str(e)}")
            return f"[Error parsing {filename}: {str(e)}]"

class LLMExtractionService:
    def __init__(self, llm_config: Dict[str, Any] = None):
        self.llm = get_llm(
            api_key=llm_config.get("api_key") if llm_config else None,
            base_url=llm_config.get("base_url") if llm_config else None,
            model_name=llm_config.get("model_id") if llm_config else None
        )
        self.chunk_size = 20000

    async def extract_and_summarize(
        self, 
        text: str, 
        concurrency: int = 3, 
        status_callback: Callable[[str], Any] = None
    ) -> AsyncGenerator[Dict[str, Any], None]:
        """Chunks text and processes them in parallel using LLM, yielding results as they complete."""
        if not text:
            return

        chunks = [text[i:i + self.chunk_size] for i in range(0, len(text), self.chunk_size)]
        total_chunks = len(chunks)
        
        if status_callback:
            res = status_callback(f"Total chunks to process: {total_chunks}")
            if asyncio.iscoroutine(res): await res

        semaphore = asyncio.Semaphore(concurrency)

        async def process_chunk(index: int, chunk: str) -> Dict[str, Any]:
            async with semaphore:
                if status_callback:
                    res = status_callback(f"Summarizing chunk {index + 1}/{total_chunks}...")
                    if asyncio.iscoroutine(res): await res
                
                system_prompt = (
                    "You are a highly skilled Data Extraction and Analysis Specialist. Your goal is to convert the provided text into a high-density information summary that will be used for diagram generation (flowcharts, mind maps, timelines, etc.).\n\n"
                    "Please extract the following elements with high precision:\n"
                    "1. **Temporal Data**: All dates, times, durations, and chronological sequences.\n"
                    "2. **Key Entities**: Names of people, organizations, systems, and specialized terms.\n"
                    "3. **Core Relationships**: How entities interact, causal links, and hierarchical dependencies.\n"
                    "4. **Quantitative Specifications**: Measurements, percentages, financial figures, and technical specs.\n"
                    "5. **Procedural Logic**: Step-by-step processes, decision points, and conditional flows.\n\n"
                    "Format your output as a structured Markdown summary that is clear, logical, and optimized for downstream AI reasoning."
                ) + get_time_instructions()
                messages = [
                    SystemMessage(content=system_prompt),
                    HumanMessage(content=f"Text chunk {index + 1}/{total_chunks}:\n\n{chunk}")
                ]
                
                try:
                    response = await self.llm.ainvoke(messages)
                    return {"index": index, "content": response.content}
                except Exception as e:
                    logger.error(f"Error processing chunk {index + 1}: {str(e)}")
                    return {"index": index, "content": f"[Error in chunk {index + 1}]"}

        tasks = [process_chunk(i, chunk) for i, chunk in enumerate(chunks)]
        
        summaries = [None] * total_chunks
        for task in asyncio.as_completed(tasks):
            result = await task
            summaries[result["index"]] = result["content"]
            yield result

        # Final synthesis if multiple chunks
        if len(chunks) > 1:
            if status_callback:
                res = status_callback("Synthesizing final summary...")
                if asyncio.iscoroutine(res): await res
            
            combined_summaries = "\n\n---\n\n".join([s for s in summaries if s])
            final_system = (
                "You are a Master Synthesis Architect. You will receive several partial summaries extracted from a larger document. Your task is to unify them into a single, cohesive, and comprehensive 'Master Intelligence Document'.\n\n"
                "Your final synthesis must:\n"
                "1. **Eliminate Redundancy**: Merge overlapping information into a crisp structure.\n"
                "2. **Enforce Chronology**: If the content involves processes or history, ensure a logical timeline.\n"
                "3. **Preserve Depth**: Do not lose specific technical details, key metrics, or critical dates.\n"
                "4. **Optimize for Visualization**: Structure the information (using nested headers, lists, and tables where appropriate) "
                "such that it can be easily transformed into architectural diagrams or logical maps.\n\n"
                "The goal is to provide the ultimate context for a diagram-generation agent to create accurate and professional visual representations of the original document."
            ) + get_time_instructions()
            final_messages = [
                SystemMessage(content=final_system),
                HumanMessage(content=f"Partial Summaries:\n\n{combined_summaries}")
            ]
            final_response = await self.llm.ainvoke(final_messages)
            yield {"index": -1, "content": final_response.content, "is_final": True}
        else:
            yield {"index": -1, "content": summaries[0] if summaries else "", "is_final": True}
